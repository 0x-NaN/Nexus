import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# Color Palette (Dark Theme / Amber Accent)
BG_COLOR = RGBColor(9, 9, 11)        # #09090b
TEXT_COLOR = RGBColor(244, 244, 245)  # #f4f4f5
ACCENT_COLOR = RGBColor(240, 180, 41) # #f0b429 (Amber)
MUTED_COLOR = RGBColor(161, 161, 170)  # #a1a1aa

def set_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_grid(slide):
    # center-left grid
    left = Inches(0.2)
    top = Inches(2.25)
    size = Inches(3.0)
    steps = 10
    step_size = size / steps
    grid_color = RGBColor(40, 32, 10) # Very dark amber
    
    for i in range(steps + 1):
        # Horizontal
        y = top + i * step_size
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, y, left + size, y)
        line.line.color.rgb = grid_color
        line.line.width = Pt(1)
        # Vertical
        x = left + i * step_size
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, top, x, top + size)
        line.line.color.rgb = grid_color
        line.line.width = Pt(1)

def parse_markdown(filepath):
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split slides by ---
    raw_slides = content.split('\n---\n')
    
    # First block is front matter
    if raw_slides[0].strip().startswith('---'):
        # In case the first split didn't catch the ending '---' properly
        parts = raw_slides[0].split('---')
        if len(parts) >= 3:
            raw_slides[0] = parts[2]
        else:
            raw_slides.pop(0)
            
    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue
        
        # Parse title
        title = ""
        title_match = re.search(r'^(?:#|##)\s+(.+)$', raw, re.MULTILINE)
        if title_match:
            title = title_match.group(1).strip()
            # Remove title line from slide content
            raw = re.sub(r'^(?:#|##)\s+.+$', '', raw, flags=re.MULTILINE).strip()
            
        # Parse image
        image_path = None
        img_match = re.search(r'!\[.*?\]\((.*?)\)', raw)
        if img_match:
            image_path = img_match.group(1).strip()
            # Remove image line from slide content
            raw = re.sub(r'!\[.*?\]\(.*?\)', '', raw).strip()
            
        # Parse annotation
        annotation = None
        anno_match = re.search(r'<div class="annotation">(.*?)</div>', raw)
        if anno_match:
            annotation = anno_match.group(1).strip()
            # Remove annotation line from slide content
            raw = re.sub(r'<div class="annotation">.*?</div>', '', raw).strip()
            
        # Parse bullet points
        bullets = []
        lines = raw.split('\n')
        remaining_lines = []
        for line in lines:
            line = line.strip()
            if not line:
                continue
            if line.startswith('- ') or line.startswith('* '):
                bullets.append(line[2:])
            else:
                remaining_lines.append(line)
                
        body_text = "\n".join(remaining_lines).strip()
        
        # Clean title markdown formatting (like strong tags)
        title = title.replace('**', '').replace('*', '')
        
        slides.append({
            'title': title,
            'image': image_path,
            'annotation': annotation,
            'bullets': bullets,
            'body': body_text
        })
        
    return slides

def build_presentation(slides, output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    
    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        set_background(slide)
        
        # Skip title slide grid for aesthetics, add to all other slides
        if i > 0:
            add_grid(slide)
            
        title_text = slide_data['title']
        bullets = slide_data['bullets']
        body_text = slide_data['body']
        image_path = slide_data['image']
        annotation = slide_data['annotation']
        
        # Title Slide (First slide)
        if i == 0:
            # Main Title
            title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(1.5))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.text = title_text if title_text else "Nexus"
            p.font.name = 'Inter'
            p.font.size = Pt(56)
            p.font.bold = True
            p.font.color.rgb = TEXT_COLOR
            
            # Subtitle
            sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(11.333), Inches(2.5))
            tf_sub = sub_box.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.alignment = PP_ALIGN.CENTER
            
            # Format the body text on title slide
            full_sub_text = body_text.replace('**', '')
            p_sub.text = full_sub_text
            p_sub.font.name = 'Inter'
            p_sub.font.size = Pt(20)
            p_sub.font.color.rgb = MUTED_COLOR
            continue
            
        # Standard Slides
        # 1. Slide Title
        if title_text:
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = 'Inter'
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = TEXT_COLOR
            
        # 2. Content Layout (Image vs Text-only)
        if image_path:
            # Left column (Text)
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.8), Inches(4.8))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            # Add body text if exists
            if body_text:
                p = tf.paragraphs[0]
                p.text = body_text.replace('**', '')
                p.font.name = 'Inter'
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_COLOR
                p.space_after = Pt(14)
                
            # Add bullets
            for j, bullet in enumerate(bullets):
                p = tf.add_paragraph() if (body_text or j > 0) else tf.paragraphs[0]
                p.text = "• " + bullet.replace('**', '')
                p.font.name = 'Inter'
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_COLOR
                p.space_after = Pt(10)
                
            # Right column (Image with Yellow Border)
            img_left = Inches(6.2)
            img_top = Inches(1.8)
            img_width = Inches(6.3)
            img_height = Inches(4.0)
            
            # Yellow border shape
            border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_left, img_top, img_width, img_height)
            border.fill.background()
            border.line.color.rgb = ACCENT_COLOR
            border.line.width = Pt(4)
            
            # Add actual image
            if os.path.exists(image_path):
                # Inset slightly inside the border
                slide.shapes.add_picture(
                    image_path, 
                    img_left + Inches(0.04), 
                    img_top + Inches(0.04), 
                    img_width - Inches(0.08), 
                    img_height - Inches(0.08)
                )
            
            # Add annotation box
            if annotation:
                anno_box = slide.shapes.add_textbox(img_left, img_top + img_height + Inches(0.1), img_width, Inches(0.8))
                tf_anno = anno_box.text_frame
                tf_anno.word_wrap = True
                p_anno = tf_anno.paragraphs[0]
                p_anno.alignment = PP_ALIGN.CENTER
                p_anno.text = annotation
                p_anno.font.name = 'Inter'
                p_anno.font.size = Pt(16)
                p_anno.font.bold = True
                p_anno.font.color.rgb = ACCENT_COLOR
                
        else:
            # Text-only Slide (takes full width)
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            # Special case for code block layouts (e.g. Solution flow, Architecture)
            if '```' in body_text:
                # Format code blocks cleanly
                code_lines = [line.strip() for line in body_text.split('\n') if not line.strip().startswith('```')]
                code_text = "\n".join(code_lines)
                
                # Title Slide or body
                p = tf.paragraphs[0]
                p.text = code_text
                p.font.name = 'JetBrains Mono'
                p.font.size = Pt(18)
                p.font.color.rgb = MUTED_COLOR
                p.space_after = Pt(14)
            else:
                if body_text:
                    p = tf.paragraphs[0]
                    p.text = body_text.replace('**', '')
                    p.font.name = 'Inter'
                    p.font.size = Pt(20)
                    p.font.color.rgb = TEXT_COLOR
                    p.space_after = Pt(14)
                
                for j, bullet in enumerate(bullets):
                    p = tf.add_paragraph() if (body_text or j > 0) else tf.paragraphs[0]
                    p.text = "• " + bullet.replace('**', '')
                    p.font.name = 'Inter'
                    p.font.size = Pt(20)
                    p.font.color.rgb = TEXT_COLOR
                    p.space_after = Pt(12)
                    
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}!")

if __name__ == '__main__':
    md_path = 'nexus-deck.md'
    pptx_path = 'nexus-deck.pptx'
    slides = parse_markdown(md_path)
    build_presentation(slides, pptx_path)
